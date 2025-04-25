# This program is licensed under the GNU General Public License v3.0.
# For more details, see: https://www.gnu.org/licenses/gpl-3.0.en.html
#
# Author: Jonathan Graziola (isidore.gpt@gmail.com)
#
# This program is free software: you can redistribute it and/or modify
# it under the terms of the GNU General Public License as published by
# the Free Software Foundation, either version 3 of the License, or
# (at your option) any later version.
#
# This program is distributed in the hope that it will be useful,
# but WITHOUT ANY WARRANTY; without even the implied warranty of
# MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE. See the
# GNU General Public License for more details.
#
# You should have received a copy of the GNU General Public License
# along with this program. If not, see <https://www.gnu.org/licenses/>.

import streamlit as st
import datetime
import io
import zipfile
import os
import tempfile
from fpdf import FPDF
import pandas as pd
from openai import OpenAI
import re
import time
import base64
from PIL import Image
import pytesseract
from pdf2image import convert_from_path, convert_from_bytes
import fitz  # PyMuPDF
import docx2txt
from pptx import Presentation
import openpyxl
import glob
import numpy as np
from io import BytesIO

# Model provider configurations
MODEL_OPTIONS = {
    "OpenAI": {
        "models": [
            "gpt-4.1-2025-04-14", "gpt-4o-2024-08-06", "gpt-4o-mini-2024-07-18"
        ],
        "supports_temperature": True,
        "temp_range": (0.0, 1.0)
    },
    "Anthropic": {
        "models": ["claude-3-5-sonnet-20241022", "claude-3-7-sonnet-20250219"],
        "supports_temperature": True,
        "temp_range": (0.0, 1.0)
    }
}

# Define document types common in private equity deals
DOCUMENT_TYPES = [
    "Confidential Information Memorandum (CIM)",
    "Executive Summary/Teaser",
    "Pitch Deck",
    "Investor Letter",
    "NDA",
    "LOI",
    "SPA",
    "Board Presentation",
    "Financial Statements",
    "Forecasts"
]

# Define common financial and legal terms to check for consistency
FINANCIAL_TERMS = [
    "EBITDA",
    "Enterprise Value",
    "Run-Rate Revenue",
    "ARR",
    "MRR",
    "Cap Table",
    "Trailing Twelve Months (TTM)",
    "Debt-to-EBITDA",
    "MOIC",
    "IRR",
    "Net Income",
    "Gross Margin",
    "Operating Margin",
    "LTM",
    "YoY",
    "QoQ"
]

class DocumentProcessor:
    """Document processing and extraction class"""
    
    @staticmethod
    def extract_from_txt(file_content):
        """Extract text from a plain text file"""
        try:
            if isinstance(file_content, bytes):
                text = file_content.decode('utf-8', errors='replace')
            else:
                text = file_content
            return {
                'text': text,
                'images': [],
                'tables': [],
                'slides': [],
                'charts': []
            }
        except Exception as e:
            return {'error': f"Error processing text file: {str(e)}"}
    
    @staticmethod
    def extract_from_docx(file_content):
        """Extract text and images from a DOCX file"""
        try:
            # Create a temp file to store the content
            with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as temp_file:
                temp_file.write(file_content)
                temp_file_path = temp_file.name
            
            # Extract text
            text = docx2txt.process(temp_file_path)
            
            # Get images
            images = []
            image_dir = tempfile.mkdtemp()
            extracted_images = docx2txt.process(temp_file_path, image_dir)
            for img_path in glob.glob(os.path.join(image_dir, "*.png")):
                with open(img_path, "rb") as img_file:
                    img_data = img_file.read()
                    images.append({
                        'data': img_data,
                        'format': 'png',
                        'filename': os.path.basename(img_path)
                    })
            
            # Clean up
            os.unlink(temp_file_path)
            
            # Return results
            return {
                'text': text,
                'images': images,
                'tables': [],  # Would need more complex parsing for tables in DOCX
                'slides': [],
                'charts': []
            }
        except Exception as e:
            return {'error': f"Error processing DOCX file: {str(e)}"}
    
    @staticmethod
    def extract_from_pdf(file_content):
        """Extract text and images from a PDF file"""
        try:
            # Create a temp file to store the content
            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as temp_file:
                temp_file.write(file_content)
                temp_file_path = temp_file.name
            
            # Open the PDF with PyMuPDF
            doc = fitz.open(temp_file_path)
            
            text = ""
            images = []
            tables = []
            
            # Extract text and images page by page
            for page_num, page in enumerate(doc):
                # Extract text
                text += page.get_text()
                
                # Extract images
                for img_index, img in enumerate(page.get_images(full=True)):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_data = base_image["image"]
                    
                    images.append({
                        'data': image_data,
                        'format': base_image["ext"],
                        'filename': f"page_{page_num + 1}_img_{img_index + 1}.{base_image['ext']}"
                    })
                
                # Attempt to identify tables (simple heuristic based on layout)
                blocks = page.get_text("blocks")
                for block in blocks:
                    # Simple heuristic: if a block has multiple lines with similar width and position,
                    # it might be a table
                    if len(block[4].split('\n')) > 3 and '  ' in block[4]:
                        tables.append({
                            'page': page_num + 1,
                            'content': block[4],
                            'bbox': block[:4]  # bounding box
                        })
            
            # Clean up
            doc.close()
            os.unlink(temp_file_path)
            
            return {
                'text': text,
                'images': images,
                'tables': tables,
                'slides': [],
                'charts': []  # More advanced analysis would be needed to identify charts
            }
        except Exception as e:
            return {'error': f"Error processing PDF file: {str(e)}"}
    
    @staticmethod
    def extract_from_pptx(file_content):
        """Extract text, images and slide info from a PPTX file"""
        try:
            # Create a temp file to store the content
            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:
                temp_file.write(file_content)
                temp_file_path = temp_file.name
            
            # Open the presentation
            presentation = Presentation(temp_file_path)
            
            text = ""
            images = []
            slides = []
            
            # Extract content slide by slide
            for slide_index, slide in enumerate(presentation.slides):
                slide_text = ""
                slide_images = []
                
                # Get text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                    
                    # Extract images
                    if shape.shape_type == 13:  # 13 is the enum value for pictures
                        image = shape.image
                        image_bytes = image.blob
                        
                        image_info = {
                            'data': image_bytes,
                            'format': image.ext,
                            'filename': f"slide_{slide_index + 1}_img_{len(slide_images) + 1}.{image.ext}"
                        }
                        
                        images.append(image_info)
                        slide_images.append(image_info)
                
                # Add to overall text
                text += slide_text
                
                # Record this slide
                slides.append({
                    'number': slide_index + 1,
                    'text': slide_text,
                    'images': [img['filename'] for img in slide_images]
                })
            
            # Clean up
            os.unlink(temp_file_path)
            
            return {
                'text': text,
                'images': images,
                'tables': [],  # More complex analysis for tables in PPTX
                'slides': slides,
                'charts': []  # Would need special handling for chart objects
            }
        except Exception as e:
            return {'error': f"Error processing PPTX file: {str(e)}"}
    
    @staticmethod
    def extract_from_xlsx(file_content):
        """Extract text and data from an XLSX file"""
        try:
            # Create a temp file to store the content
            with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as temp_file:
                temp_file.write(file_content)
                temp_file_path = temp_file.name
            
            # Open the workbook
            workbook = openpyxl.load_workbook(temp_file_path, data_only=True)
            
            text = ""
            tables = []
            charts = []
            
            # Process each worksheet
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                
                # Extract text and create a text representation of the sheet
                sheet_text = f"Sheet: {sheet_name}\n"
                sheet_data = []
                
                # Determine the data range
                max_row = sheet.max_row
                max_col = sheet.max_column
                
                # Extract data as a table
                for row in range(1, min(max_row + 1, 100)):  # Limit to first 100 rows to prevent huge extractions
                    row_data = []
                    for col in range(1, min(max_col + 1, 20)):  # Limit to first 20 columns
                        cell = sheet.cell(row=row, column=col)
                        row_data.append(str(cell.value if cell.value is not None else ''))
                    
                    sheet_data.append(row_data)
                    sheet_text += ' | '.join(row_data) + '\n'
                
                # Add to the overall text content
                text += sheet_text + '\n\n'
                
                # Record the table data
                tables.append({
                    'sheet': sheet_name,
                    'data': sheet_data
                })
                
                # Extract charts (just the references, not the actual chart data)
                for chart_rel in sheet._charts:
                    charts.append({
                        'sheet': sheet_name,
                        'type': 'Excel chart'
                    })
            
            # Clean up
            os.unlink(temp_file_path)
            
            return {
                'text': text,
                'images': [],  # Excel files may have images but extra processing needed
                'tables': tables,
                'slides': [],
                'charts': charts
            }
        except Exception as e:
            return {'error': f"Error processing XLSX file: {str(e)}"}
    
    @staticmethod
    def analyze_image(image_data, image_format):
        """Perform OCR and basic analysis on an image"""
        try:
            # Create a PIL Image from the binary data
            img = Image.open(BytesIO(image_data))
            
            # Perform OCR to extract text
            text = pytesseract.image_to_string(img)
            
            # Get image dimensions and format
            width, height = img.size
            
            # Basic image quality assessment (example)
            is_small = width < 300 or height < 300
            
            return {
                'text': text,
                'width': width,
                'height': height,
                'format': image_format,
                'is_small': is_small
            }
        except Exception as e:
            return {'error': f"Error analyzing image: {str(e)}"}
    
    @staticmethod
    def process_document(file_content, file_name):
        """Process a document based on its file extension"""
        file_extension = os.path.splitext(file_name)[1].lower()
        
        if file_extension == '.txt' or file_extension == '.md':
            return DocumentProcessor.extract_from_txt(file_content)
        elif file_extension == '.docx':
            return DocumentProcessor.extract_from_docx(file_content)
        elif file_extension == '.pdf':
            return DocumentProcessor.extract_from_pdf(file_content)
        elif file_extension == '.pptx':
            return DocumentProcessor.extract_from_pptx(file_content)
        elif file_extension == '.xlsx':
            return DocumentProcessor.extract_from_xlsx(file_content)
        else:
            # Try to process as text for unknown formats
            try:
                text = file_content.decode('utf-8', errors='replace')
                return {
                    'text': text,
                    'images': [],
                    'tables': [],
                    'slides': [],
                    'charts': []
                }
            except:
                return {'error': f"Unsupported file format: {file_extension}"}


class ProofreadingManager:
    def __init__(self, api_key):
        self.api_key = api_key
        self.openai_client = OpenAI(api_key=api_key)
        
    def proofread_document_with_openai(self, model, document_data, document_type, style_guide=None, temperature=0.1):
        """Proofread a document using an OpenAI model"""
        try:
            # Extract main text content
            document_text = document_data.get('text', '')
            
            # Count visual elements
            num_images = len(document_data.get('images', []))
            num_tables = len(document_data.get('tables', []))
            num_slides = len(document_data.get('slides', []))
            num_charts = len(document_data.get('charts', []))
            
            # Prepare the proofreading system prompt
            system_prompt = """You are an expert proofreader for private equity deal documents. Your task is to:
1. Correct grammar, spelling, and punctuation errors
2. Ensure consistency in terminology (especially financial terms like EBITDA, enterprise value, etc.)
3. Verify document formatting and style is uniform 
4. Maintain professional tone and clarity
5. Identify ambiguities or confusing phrasing that could create legal or financial misunderstandings
6. Check that dates, milestones, and party names are used consistently
7. Flag potential compliance or legal inconsistencies (at a surface level)

Provide your feedback in these sections:
- SUMMARY: Brief overview of document quality
- CRITICAL ISSUES: High-priority items that must be fixed
- CONSISTENCY CONCERNS: Terminology or style inconsistencies
- GRAMMAR & SPELLING: List of grammatical or spelling errors
- FORMATTING: Issues with document layout and structure
- SUGGESTIONS: Recommendations for clarity and professionalism
- REVISED TEXT: The corrected version of the document with all fixes implemented

For documents with visual elements (images, charts, tables, slides), also include:
- VISUAL ELEMENTS: Assessment of how well the visual elements support the document's purpose
- LAYOUT CONCERNS: Issues with the positioning or formatting of visual elements

When identifying errors, please use the following format when possible:
- Original: "text with error" should be "corrected text"

This format will help automatically identify and track corrections.

Do not disclose confidential information in your feedback. Focus on improving the quality while maintaining the original meaning.
"""

            # Add style guide information if provided
            if style_guide:
                system_prompt += f"\n\nApply the following style guide requirements:\n{style_guide}"
                
            # Prepare the proofreading request with info about visual elements
            proofreading_request = f"""Document Type: {document_type}

Document Information:
- Contains {num_images} images
- Contains {num_tables} tables
- Contains {num_slides} slides
- Contains {num_charts} charts

Please proofread the following document:

{document_text}
"""
            
            # If there are images with extracted text, include that info too
            if num_images > 0 and 'images' in document_data:
                image_text = "\n\nText extracted from images:\n"
                for i, img in enumerate(document_data.get('images', [])[:5]):  # Limit to first 5 images
                    if 'data' in img and img['data']:
                        # Try to analyze the image and extract text
                        analysis = DocumentProcessor.analyze_image(img['data'], img.get('format', 'unknown'))
                        if 'text' in analysis and analysis['text'].strip():
                            image_text += f"\nImage {i+1} ({img.get('filename', 'unnamed')}): {analysis['text'].strip()}\n"
                
                proofreading_request += image_text
            
            # Create chat completion request
            response = self.openai_client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": proofreading_request}
                ],
                temperature=temperature
            )
            
            return {
                "content": response.choices[0].message.content
            }
            
        except Exception as e:
            return {
                "content": f"Error proofreading with OpenAI model {model}: {str(e)}"
            }
    
    def proofread_document_with_anthropic(self, model, document_data, document_type, style_guide=None, temperature=0.1):
        """Proofread a document using an Anthropic model"""
        try:
            import requests
            
            # Extract main text content
            document_text = document_data.get('text', '')
            
            # Count visual elements
            num_images = len(document_data.get('images', []))
            num_tables = len(document_data.get('tables', []))
            num_slides = len(document_data.get('slides', []))
            num_charts = len(document_data.get('charts', []))
            
            # Prepare the proofreading system prompt
            system_prompt = """You are an expert proofreader for private equity deal documents. Your task is to:
1. Correct grammar, spelling, and punctuation errors
2. Ensure consistency in terminology (especially financial terms like EBITDA, enterprise value, etc.)
3. Verify document formatting and style is uniform 
4. Maintain professional tone and clarity
5. Identify ambiguities or confusing phrasing that could create legal or financial misunderstandings
6. Check that dates, milestones, and party names are used consistently
7. Flag potential compliance or legal inconsistencies (at a surface level)

Provide your feedback in these sections:
- SUMMARY: Brief overview of document quality
- CRITICAL ISSUES: High-priority items that must be fixed
- CONSISTENCY CONCERNS: Terminology or style inconsistencies
- GRAMMAR & SPELLING: List of grammatical or spelling errors
- FORMATTING: Issues with document layout and structure
- SUGGESTIONS: Recommendations for clarity and professionalism
- REVISED TEXT: The corrected version of the document with all fixes implemented

For documents with visual elements (images, charts, tables, slides), also include:
- VISUAL ELEMENTS: Assessment of how well the visual elements support the document's purpose
- LAYOUT CONCERNS: Issues with the positioning or formatting of visual elements

When identifying errors, please use the following format when possible:
- Original: "text with error" should be "corrected text"

This format will help automatically identify and track corrections.

Do not disclose confidential information in your feedback. Focus on improving the quality while maintaining the original meaning.
"""

            # Add style guide information if provided
            if style_guide:
                system_prompt += f"\n\nApply the following style guide requirements:\n{style_guide}"
                
            # Prepare the proofreading request with info about visual elements
            proofreading_request = f"""Document Type: {document_type}

Document Information:
- Contains {num_images} images
- Contains {num_tables} tables
- Contains {num_slides} slides
- Contains {num_charts} charts

Please proofread the following document:

{document_text}
"""
            
            # If there are images with extracted text, include that info too
            if num_images > 0 and 'images' in document_data:
                image_text = "\n\nText extracted from images:\n"
                for i, img in enumerate(document_data.get('images', [])[:5]):  # Limit to first 5 images
                    if 'data' in img and img['data']:
                        # Try to analyze the image and extract text
                        analysis = DocumentProcessor.analyze_image(img['data'], img.get('format', 'unknown'))
                        if 'text' in analysis and analysis['text'].strip():
                            image_text += f"\nImage {i+1} ({img.get('filename', 'unnamed')}): {analysis['text'].strip()}\n"
                
                proofreading_request += image_text
            
            headers = {
                "Content-Type": "application/json",
                "x-api-key": self.api_key,
                "anthropic-version": "2023-06-01"
            }
            
            data = {
                "model": model,
                "system": system_prompt,
                "messages": [
                    {"role": "user", "content": proofreading_request}
                ],
                "temperature": temperature,
                "max_tokens": 4000
            }
            
            response = requests.post("https://api.anthropic.com/v1/messages", json=data, headers=headers)
            response.raise_for_status()
            
            return {
                "content": response.json()["content"][0]["text"]
            }
            
        except Exception as e:
            return {
                "content": f"Error proofreading with Anthropic model {model}: {str(e)}"
            }
    
    def proofread_document(self, provider, model, document_data, document_type, style_guide=None, temperature=0.1):
        """Route the proofreading request to the appropriate API based on provider"""
        if provider == "OpenAI":
            return self.proofread_document_with_openai(model, document_data, document_type, style_guide, temperature)
        elif provider == "Anthropic":
            return self.proofread_document_with_anthropic(model, document_data, document_type, style_guide, temperature)
        else:
            return {
                "content": f"Unsupported provider: {provider}"
            }

    def extract_financial_terms(self, document_text):
        """Extract and analyze financial terms usage in the document"""
        terms_count = {}
        variants = {}
        
        for term in FINANCIAL_TERMS:
            # Look for the term and potential variants
            base_term = term.split(" ")[0]  # Take first word if multi-word term
            pattern = rf'\b{re.escape(base_term)}[-\s]*(?:\([^)]*\))?[-\s]*\w*\b'
            matches = re.findall(pattern, document_text, re.IGNORECASE)
            
            if matches:
                terms_count[term] = len(matches)
                
                # Get unique variants
                unique_variants = list(set([m.strip() for m in matches]))
                if len(unique_variants) > 1:
                    variants[term] = unique_variants
        
        return {
            "terms_count": terms_count,
            "variants": variants
        }
    
    def check_monetary_formatting(self, document_text):
        """Check consistency of monetary value formatting"""
        # Look for different formats like $25 million, $25M, $25MM, etc.
        patterns = {
            'spelled_out': r'\$\s*\d+(?:\.\d+)?\s*(?:million|billion|thousand)',
            'M_format': r'\$\s*\d+(?:\.\d+)?\s*M\b',
            'MM_format': r'\$\s*\d+(?:\.\d+)?\s*MM\b',
            'K_format': r'\$\s*\d+(?:\.\d+)?\s*K\b',
            'B_format': r'\$\s*\d+(?:\.\d+)?\s*B\b'
        }
        
        results = {}
        for name, pattern in patterns.items():
            matches = re.findall(pattern, document_text)
            if matches:
                results[name] = matches
        
        # Determine if there's inconsistency in formatting
        inconsistent = len(results) > 1
        
        return {
            "formats": results,
            "inconsistent": inconsistent
        }
    
    def analyze_dates(self, document_text):
        """Extract and analyze date formats in the document"""
        # Common date formats
        date_patterns = {
            'mm/dd/yyyy': r'\b\d{1,2}/\d{1,2}/\d{4}\b',
            'dd/mm/yyyy': r'\b\d{1,2}/\d{1,2}/\d{4}\b',  # Note: overlaps with mm/dd/yyyy
            'month_dd_yyyy': r'\b(?:January|February|March|April|May|June|July|August|September|October|November|December)\s+\d{1,2},\s+\d{4}\b',
            'dd_month_yyyy': r'\b\d{1,2}\s+(?:January|February|March|April|May|June|July|August|September|October|November|December)\s+\d{4}\b',
            'yyyy-mm-dd': r'\b\d{4}-\d{1,2}-\d{1,2}\b'
        }
        
        results = {}
        for name, pattern in date_patterns.items():
            matches = re.findall(pattern, document_text)
            if matches:
                results[name] = matches
        
        # Determine if there's inconsistency in date formatting
        inconsistent = len(results) > 1
        
        return {
            "formats": results,
            "inconsistent": inconsistent
        }


def generate_error_summary(document_text, ai_results):
    """Generate a structured error summary from AI proofreading results"""
    import re
    import pandas as pd
    import io
    
    all_errors = []
    
    # Process each AI model's feedback
    for model, result in ai_results.items():
        content = result["content"]
        
        # Attempt to extract the errors section
        grammar_section = re.search(r'GRAMMAR & SPELLING:(.*?)(?:FORMATTING:|SUGGESTIONS:|REVISED TEXT:|VISUAL ELEMENTS:|$)', 
                                   content, re.DOTALL | re.IGNORECASE)
        consistency_section = re.search(r'CONSISTENCY CONCERNS:(.*?)(?:GRAMMAR & SPELLING:|FORMATTING:|SUGGESTIONS:|REVISED TEXT:|VISUAL ELEMENTS:|$)', 
                                      content, re.DOTALL | re.IGNORECASE)
        formatting_section = re.search(r'FORMATTING:(.*?)(?:SUGGESTIONS:|REVISED TEXT:|VISUAL ELEMENTS:|$)', 
                                     content, re.DOTALL | re.IGNORECASE)
        
        # Get the revised text if available
        revised_text_match = re.search(r'REVISED TEXT:(.*?)(?:VISUAL ELEMENTS:|$)', content, re.DOTALL | re.IGNORECASE)
        revised_text = revised_text_match.group(1).strip() if revised_text_match else None
        
        # Process each type of error
        sections = [
            ("Grammar/Spelling", grammar_section.group(1).strip() if grammar_section else None),
            ("Consistency", consistency_section.group(1).strip() if consistency_section else None),
            ("Formatting", formatting_section.group(1).strip() if formatting_section else None)
        ]
        
        for error_type, section_text in sections:
            if not section_text:
                continue
                
            # Try to extract individual errors - this is a bit heuristic
            # Look for bullet points, numbered lists, or just newlines
            error_items = re.findall(r'[-•*]\s+(.*?)(?=[-•*]|\n\n|$)', section_text, re.DOTALL)
            if not error_items:
                error_items = re.findall(r'\d+\.\s+(.*?)(?=\d+\.|\n\n|$)', section_text, re.DOTALL)
            if not error_items:
                error_items = section_text.split('\n')
                
            for error_item in error_items:
                error_item = error_item.strip()
                if not error_item:
                    continue
                    
                # Try to find the original text and corrected version
                # Look for patterns like "X should be Y" or "Change X to Y"
                correction_match = re.search(r'"([^"]+)"\s+(?:should be|change to|replace with|correct to)\s+"([^"]+)"', 
                                           error_item, re.IGNORECASE)
                
                # Also look for "Original: X should be Y" format
                if not correction_match:
                    correction_match = re.search(r'Original:\s*"([^"]+)"\s+(?:should be|change to|replace with|correct to)\s+"([^"]+)"', 
                                               error_item, re.IGNORECASE)
                
                if correction_match:
                    original = correction_match.group(1).strip()
                    corrected = correction_match.group(2).strip()
                    
                    # Find location in document
                    location = "Unknown"
                    try:
                        # If original text is found in document, get position
                        if original in document_text:
                            pos = document_text.find(original)
                            # Get surrounding context to help locate
                            start = max(0, pos - 20)
                            end = min(len(document_text), pos + len(original) + 20)
                            context = document_text[start:end]
                            
                            # Count line number
                            line_num = document_text[:pos].count('\n') + 1
                            location = f"Line {line_num} (approx.): ...{context}..."
                    except:
                        pass
                    
                    all_errors.append({
                        "Model": model,
                        "Error Type": error_type,
                        "Original Text": original,
                        "Corrected Text": corrected,
                        "Location": location,
                        "Description": error_item
                    })
                else:
                    # If we can't parse the correction clearly, just include the full error description
                    all_errors.append({
                        "Model": model,
                        "Error Type": error_type,
                        "Original Text": "",
                        "Corrected Text": "",
                        "Location": "N/A",
                        "Description": error_item
                    })
    
    # Create a DataFrame from all errors
    df = pd.DataFrame(all_errors)
    
    # Generate CSV and text summaries
    csv_buffer = io.StringIO()
    df.to_csv(csv_buffer, index=False)
    csv_text = csv_buffer.getvalue()
    

    # More human-readable text format
    text_summary = "ERROR SUMMARY REPORT\n\n"
    text_summary += f"Total Errors Found: {len(all_errors)}\n\n"
    
    for model in set(df["Model"]):
        model_errors = df[df["Model"] == model]
        text_summary += f"===== {model} =====\n"
        text_summary += f"Found {len(model_errors)} errors\n\n"
        
        for idx, error in model_errors.iterrows():
            text_summary += f"ERROR #{idx+1}: {error['Error Type']}\n"
            text_summary += f"Description: {error['Description']}\n"
            if error['Original Text']:
                text_summary += f"Original: \"{error['Original Text']}\"\n"
            if error['Corrected Text']:
                text_summary += f"Corrected: \"{error['Corrected Text']}\"\n"
            if error['Location'] != "N/A":
                text_summary += f"Location: {error['Location']}\n"
            text_summary += "\n"
    
    return {
        "csv": csv_text.encode('utf-8'),
        "text": text_summary.encode('utf-8'),
        "dataframe": df
    }
